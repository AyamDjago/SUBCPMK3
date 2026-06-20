"""
Generate presentasi_cbr.pptx - 5 Slide Presentasi Tugas Akhir CBR
Universitas Muhammadiyah Malang - Penalaran Komputer
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

# ==============================================================
# WARNA TEMA (UMM Merah + Dark Profesional)
# ==============================================================
C_BG_DARK    = RGBColor(0x14, 0x18, 0x2B)   # Navy very dark
C_ACCENT     = RGBColor(0xC0, 0x1C, 0x2E)   # Merah UMM
C_ACCENT2    = RGBColor(0xE8, 0x3A, 0x4D)   # Merah terang
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xF0)
C_GOLD       = RGBColor(0xF5, 0xC5, 0x18)   # Gold accent
C_PANEL      = RGBColor(0x1E, 0x24, 0x3C)   # Panel dark
C_TEXT_DIM   = RGBColor(0xA8, 0xB0, 0xCC)   # Subdued text

# WAJIB: slide 16:9
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# Helper shortcut
W = prs.slide_width
H = prs.slide_height

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, italic=False,
                color=C_WHITE, align=PP_ALIGN.LEFT,
                word_wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_multi_line_textbox(slide, lines, left, top, width, height,
                            font_size=14, color=C_WHITE, font_name="Calibri",
                            bold=False, line_spacing=1.2, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox

# ==============================================================
# SLIDE 1: JUDUL
# ==============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, C_BG_DARK)

# Top accent bar
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), fill_color=C_ACCENT)

# Left vertical accent bar
add_rect(slide, Inches(0), Inches(0.08), Inches(0.15), H - Inches(0.08), fill_color=C_PANEL)

# Decorative right panel
add_rect(slide, W - Inches(3.2), Inches(0.08), Inches(3.2), H - Inches(0.08), fill_color=C_PANEL)

# Red vertical stripe in right panel
add_rect(slide, W - Inches(3.2), Inches(0.08), Inches(0.06), H - Inches(0.08), fill_color=C_ACCENT)

# Gold accent line under title area
add_rect(slide, Inches(0.5), Inches(4.8), Inches(6.2), Inches(0.04), fill_color=C_GOLD)

# JUDUL
add_textbox(slide,
    "Penerapan Case-Based Reasoning\nBerbasis IndoBERT untuk Analisis\nPutusan Pidana Pemalsuan",
    left=Inches(0.5), top=Inches(1.2), width=Inches(8.8), height=Inches(2.8),
    font_size=30, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT, font_name="Calibri"
)

# Sub-label
add_textbox(slide,
    "Tugas Akhir · Mata Kuliah Penalaran Komputer",
    left=Inches(0.5), top=Inches(4.0), width=Inches(7.0), height=Inches(0.5),
    font_size=13, color=C_ACCENT2, align=PP_ALIGN.LEFT, italic=True
)

# Author info in right panel
add_multi_line_textbox(slide,
    ["[NAMA LENGKAP ANDA]", "", "NIM: [NIM ANDA]", "",
     "Teknik Informatika", "Universitas Muhammadiyah", "Malang", "",
     "2025 / 2026"],
    left=W - Inches(3.0), top=Inches(1.5), width=Inches(2.8), height=Inches(5.5),
    font_size=13, color=C_TEXT_DIM, align=PP_ALIGN.CENTER
)

# UMM label
add_textbox(slide, "UNIVERSITAS MUHAMMADIYAH MALANG",
    left=Inches(0.5), top=Inches(5.1), width=Inches(7.0), height=Inches(0.4),
    font_size=11, color=C_TEXT_DIM, align=PP_ALIGN.LEFT)

# Slide number
add_textbox(slide, "01 / 05",
    left=W - Inches(1.3), top=H - Inches(0.4), width=Inches(1.0), height=Inches(0.3),
    font_size=10, color=C_TEXT_DIM, align=PP_ALIGN.RIGHT)

# Bottom bar
add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=C_ACCENT)

# ==============================================================
# SLIDE 2: METODOLOGI CBR
# ==============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), fill_color=C_ACCENT)
add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=C_ACCENT)

# Header bar
add_rect(slide, Inches(0), Inches(0.08), W, Inches(0.85), fill_color=C_PANEL)
add_textbox(slide, "METODOLOGI — Siklus Case-Based Reasoning",
    left=Inches(0.3), top=Inches(0.15), width=W - Inches(0.6), height=Inches(0.6),
    font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

# Dataset info box
add_rect(slide, Inches(0.3), Inches(1.1), Inches(3.0), Inches(1.8), fill_color=C_PANEL)
add_rect(slide, Inches(0.3), Inches(1.1), Inches(0.05), Inches(1.8), fill_color=C_ACCENT)
add_multi_line_textbox(slide,
    ["📂 DATASET", "35 Putusan Pidana Pemalsuan", "Mahkamah Agung RI",
     "Pasal 263–276 KUHP", "Tahun 2008–2025"],
    left=Inches(0.5), top=Inches(1.15), width=Inches(2.7), height=Inches(1.7),
    font_size=12, color=C_LIGHT_GRAY)

# 5 CBR stages
stages = [
    ("1\nCase Base", "35 PDF → TXT\nvia pdfminer.six\n7-step cleaning"),
    ("2\nRepresentasi", "18 atribut regex\nno_perkara, vonis\nringkasan_fakta"),
    ("3\nRetrieval", "IndoBERT embed.\n768-dim\nCosine Similarity"),
    ("4\nReuse", "Weighted Vote\nTop-K=5\nPrediksi vonis"),
    ("5\nEvaluasi", "Precision@K\nMRR, Accuracy\nF1-Score"),
]
stage_colors = [C_ACCENT, RGBColor(0x8B,0x14,0x28), RGBColor(0x6B,0x10,0x20),
                RGBColor(0x4A,0x0C,0x18), RGBColor(0x2A,0x08,0x10)]

box_w = Inches(2.1)
box_h = Inches(2.8)
gap = Inches(0.22)
start_x = Inches(0.3)
start_y = Inches(3.0)

for i, (label, detail) in enumerate(stages):
    x = start_x + i * (box_w + gap)
    add_rect(slide, x, start_y, box_w, box_h, fill_color=C_PANEL)
    add_rect(slide, x, start_y, box_w, Inches(0.06), fill_color=stage_colors[i] if i < len(stage_colors) else C_ACCENT)

    # Number + title
    num, title = label.split("\n", 1)
    add_textbox(slide, num,
        left=x + Inches(0.1), top=start_y + Inches(0.12), width=box_w - Inches(0.2), height=Inches(0.5),
        font_size=22, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)
    add_textbox(slide, title,
        left=x + Inches(0.1), top=start_y + Inches(0.6), width=box_w - Inches(0.2), height=Inches(0.5),
        font_size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_multi_line_textbox(slide, detail.split("\n"),
        left=x + Inches(0.1), top=start_y + Inches(1.1), width=box_w - Inches(0.2), height=Inches(1.5),
        font_size=11, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)

    # Arrow between boxes (except last)
    if i < len(stages) - 1:
        arrow_x = x + box_w + Inches(0.04)
        add_textbox(slide, "→",
            left=arrow_x, top=start_y + Inches(1.2), width=gap - Inches(0.04), height=Inches(0.4),
            font_size=18, color=C_ACCENT, align=PP_ALIGN.CENTER, bold=True)

add_textbox(slide, "02 / 05",
    left=W - Inches(1.3), top=H - Inches(0.4), width=Inches(1.0), height=Inches(0.3),
    font_size=10, color=C_TEXT_DIM, align=PP_ALIGN.RIGHT)

# ==============================================================
# SLIDE 3: IMPLEMENTASI & MODEL
# ==============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), fill_color=C_ACCENT)
add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=C_ACCENT)
add_rect(slide, Inches(0), Inches(0.08), W, Inches(0.85), fill_color=C_PANEL)
add_textbox(slide, "IMPLEMENTASI — IndoBERT & Pipeline",
    left=Inches(0.3), top=Inches(0.15), width=W - Inches(0.6), height=Inches(0.6),
    font_size=20, bold=True, color=C_WHITE)

# Left column: Model specs
add_rect(slide, Inches(0.3), Inches(1.1), Inches(5.8), Inches(5.6), fill_color=C_PANEL)
add_rect(slide, Inches(0.3), Inches(1.1), Inches(0.05), Inches(5.6), fill_color=C_ACCENT)
add_textbox(slide, "🤖  Model IndoBERT",
    left=Inches(0.5), top=Inches(1.2), width=Inches(5.5), height=Inches(0.5),
    font_size=15, bold=True, color=C_ACCENT2)

model_specs = [
    ("Model ID",        "indobenchmark/indobert-base-p1"),
    ("Arsitektur",      "BERT-Base (12 layers, 12 heads)"),
    ("Jumlah Parameter","124.441.344"),
    ("Dimensi Embedding","768"),
    ("Max Token Length", "512"),
    ("Pooling Strategy", "Mean Pooling"),
    ("Normalisasi",     "L2 Normalization"),
    ("Batch Size",      "8"),
    ("Perangkat",       "CPU"),
]
for idx, (k, v) in enumerate(model_specs):
    y = Inches(1.75) + idx * Inches(0.48)
    add_textbox(slide, k + ":",
        left=Inches(0.5), top=y, width=Inches(2.2), height=Inches(0.4),
        font_size=11, bold=True, color=C_TEXT_DIM)
    add_textbox(slide, v,
        left=Inches(2.75), top=y, width=Inches(3.2), height=Inches(0.4),
        font_size=11, color=C_WHITE)

# Right column: Pipeline steps
add_rect(slide, Inches(6.5), Inches(1.1), Inches(6.5), Inches(5.6), fill_color=C_PANEL)
add_rect(slide, Inches(6.5), Inches(1.1), Inches(0.05), Inches(5.6), fill_color=C_GOLD)
add_textbox(slide, "⚙️  Pipeline Notebook",
    left=Inches(6.7), top=Inches(1.2), width=Inches(6.1), height=Inches(0.5),
    font_size=15, bold=True, color=C_GOLD)

pipeline_steps = [
    ("01_case_base.ipynb", "35 PDF → TXT (pdfminer.six)\n7-langkah cleaning pipeline"),
    ("02_case_representation.ipynb", "Ekstraksi 18 atribut\ndengan pola regex"),
    ("03_case_retrieval.ipynb", "Generate embedding IndoBERT\nSplit 80:20 (28 train / 7 test)"),
    ("04_solution_reuse.ipynb", "Weighted Similarity Voting\nTop-K=5 prediksi vonis"),
    ("05_evaluation.ipynb", "Hitung Precision@K, MRR\nAccuracy, F1, plot grafik"),
]
for idx, (nb, desc) in enumerate(pipeline_steps):
    y = Inches(1.75) + idx * Inches(0.95)
    # Step circle/badge
    add_rect(slide, Inches(6.7), y + Inches(0.05), Inches(0.35), Inches(0.35),
             fill_color=C_ACCENT)
    add_textbox(slide, str(idx+1),
        left=Inches(6.7), top=y + Inches(0.05), width=Inches(0.35), height=Inches(0.35),
        font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, nb,
        left=Inches(7.15), top=y, width=Inches(5.7), height=Inches(0.35),
        font_size=11, bold=True, color=C_ACCENT2)
    add_textbox(slide, desc,
        left=Inches(7.15), top=y + Inches(0.35), width=Inches(5.7), height=Inches(0.5),
        font_size=10, color=C_TEXT_DIM)

add_textbox(slide, "03 / 05",
    left=W - Inches(1.3), top=H - Inches(0.4), width=Inches(1.0), height=Inches(0.3),
    font_size=10, color=C_TEXT_DIM, align=PP_ALIGN.RIGHT)

# ==============================================================
# SLIDE 4: HASIL EVALUASI
# ==============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), fill_color=C_ACCENT)
add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=C_ACCENT)
add_rect(slide, Inches(0), Inches(0.08), W, Inches(0.85), fill_color=C_PANEL)
add_textbox(slide, "HASIL & EVALUASI — Performa Sistem CBR",
    left=Inches(0.3), top=Inches(0.15), width=W - Inches(0.6), height=Inches(0.6),
    font_size=20, bold=True, color=C_WHITE)

# Key metric cards (top row)
metrics_top = [
    ("71.43%", "Accuracy"),
    ("51.02%", "Precision\n(weighted)"),
    ("71.43%", "Recall\n(weighted)"),
    ("59.52%", "F1-Score\n(weighted)"),
]
card_w = Inches(2.8)
card_h = Inches(1.5)
for i, (val, label) in enumerate(metrics_top):
    cx = Inches(0.3) + i * (card_w + Inches(0.2))
    add_rect(slide, cx, Inches(1.1), card_w, card_h, fill_color=C_PANEL)
    add_rect(slide, cx, Inches(1.1), card_w, Inches(0.05), fill_color=C_ACCENT)
    add_textbox(slide, val,
        left=cx, top=Inches(1.18), width=card_w, height=Inches(0.8),
        font_size=28, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, label,
        left=cx, top=Inches(1.95), width=card_w, height=Inches(0.6),
        font_size=11, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)

# Retrieval metrics table (bottom left)
add_rect(slide, Inches(0.3), Inches(2.8), Inches(5.8), Inches(3.9), fill_color=C_PANEL)
add_rect(slide, Inches(0.3), Inches(2.8), Inches(0.05), Inches(3.9), fill_color=C_GOLD)
add_textbox(slide, "📊  Metrik Retrieval",
    left=Inches(0.5), top=Inches(2.85), width=Inches(5.5), height=Inches(0.4),
    font_size=13, bold=True, color=C_GOLD)

ret_headers = ["K", "Precision@K", "Recall@K", "F1@K", "MRR"]
ret_data = [
    ["1", "0.0000", "0.0000", "0.0000", "0.0000"],
    ["3", "0.0000", "0.0000", "0.0000", "0.0000"],
    ["5", "0.0667", "0.1111", "0.0833", "0.0667"],
]
col_w = Inches(1.06)
row_h = Inches(0.48)
hdr_top = Inches(3.28)
for ci, hdr in enumerate(ret_headers):
    hx = Inches(0.45) + ci * col_w
    add_rect(slide, hx, hdr_top, col_w - Inches(0.04), row_h, fill_color=RGBColor(0x2A,0x30,0x50))
    add_textbox(slide, hdr, left=hx, top=hdr_top + Inches(0.05), width=col_w - Inches(0.04), height=row_h - Inches(0.1),
                font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
for ri, row_d in enumerate(ret_data):
    for ci, cell in enumerate(row_d):
        cx = Inches(0.45) + ci * col_w
        cy = hdr_top + (ri+1) * row_h
        bg = C_PANEL if ri % 2 == 0 else RGBColor(0x1A,0x20,0x38)
        add_rect(slide, cx, cy, col_w - Inches(0.04), row_h, fill_color=bg)
        col = C_ACCENT2 if (ri == 2 and ci > 0) else C_WHITE
        add_textbox(slide, cell, left=cx, top=cy + Inches(0.05), width=col_w - Inches(0.04), height=row_h - Inches(0.1),
                    font_size=10, color=col, align=PP_ALIGN.CENTER)

# Classification report (bottom right)
add_rect(slide, Inches(6.5), Inches(2.8), Inches(6.5), Inches(3.9), fill_color=C_PANEL)
add_rect(slide, Inches(6.5), Inches(2.8), Inches(0.05), Inches(3.9), fill_color=C_ACCENT)
add_textbox(slide, "📋  Classification Report",
    left=Inches(6.7), top=Inches(2.85), width=Inches(6.0), height=Inches(0.4),
    font_size=13, bold=True, color=C_ACCENT2)

cr_headers = ["Kelas", "Precision", "Recall", "F1-Score", "Support"]
cr_data = [
    ["BEBAS", "0.00", "0.00", "0.00", "2"],
    ["PENJARA_LAIN", "0.71", "1.00", "0.83", "5"],
    ["Macro avg", "0.36", "0.50", "0.42", "7"],
    ["Weighted avg", "0.51", "0.71", "0.60", "7"],
]
cr_col_w = Inches(1.18)
cr_hdr_top = Inches(3.28)
for ci, hdr in enumerate(cr_headers):
    hx = Inches(6.7) + ci * cr_col_w
    add_rect(slide, hx, cr_hdr_top, cr_col_w - Inches(0.04), row_h, fill_color=RGBColor(0x2A,0x30,0x50))
    add_textbox(slide, hdr, left=hx, top=cr_hdr_top + Inches(0.05), width=cr_col_w - Inches(0.04), height=row_h - Inches(0.1),
                font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
for ri, row_d in enumerate(cr_data):
    for ci, cell in enumerate(row_d):
        cx = Inches(6.7) + ci * cr_col_w
        cy = cr_hdr_top + (ri+1) * row_h
        bg = C_PANEL if ri % 2 == 0 else RGBColor(0x1A,0x20,0x38)
        add_rect(slide, cx, cy, cr_col_w - Inches(0.04), row_h, fill_color=bg)
        col = C_GOLD if (ri == 1 and ci in [1,2,3]) else C_WHITE
        if ri == 0 and ci in [1,2,3]:
            col = C_ACCENT2  # Red for BEBAS zero metrics
        add_textbox(slide, cell, left=cx, top=cy + Inches(0.05), width=cr_col_w - Inches(0.04), height=row_h - Inches(0.1),
                    font_size=10, color=col, align=PP_ALIGN.CENTER)

# Try to insert evaluation plot
eval_plot_path = os.path.join(os.path.dirname(__file__), "data", "eval", "evaluation_plot.png")
if os.path.exists(eval_plot_path):
    try:
        pic = slide.shapes.add_picture(eval_plot_path,
            left=Inches(6.5), top=Inches(4.9), width=Inches(6.5), height=Inches(2.45))
        print(f"[OK] evaluation_plot.png disisipkan di Slide 4")
    except Exception as e:
        print(f"[WARN] Gagal sisipkan gambar: {e}")
else:
    add_textbox(slide, "[Gambar: evaluation_plot.png — tempatkan dari data/eval/]",
        left=Inches(6.5), top=Inches(5.2), width=Inches(6.5), height=Inches(1.5),
        font_size=10, color=C_TEXT_DIM, italic=True)

add_textbox(slide, "04 / 05",
    left=W - Inches(1.3), top=H - Inches(0.4), width=Inches(1.0), height=Inches(0.3),
    font_size=10, color=C_TEXT_DIM, align=PP_ALIGN.RIGHT)

# ==============================================================
# SLIDE 5: DISKUSI & KESIMPULAN
# ==============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, C_BG_DARK)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), fill_color=C_ACCENT)
add_rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), fill_color=C_ACCENT)
add_rect(slide, Inches(0), Inches(0.08), W, Inches(0.85), fill_color=C_PANEL)
add_textbox(slide, "DISKUSI & KESIMPULAN",
    left=Inches(0.3), top=Inches(0.15), width=W - Inches(0.6), height=Inches(0.6),
    font_size=20, bold=True, color=C_WHITE)

# Left column: limitations & failure analysis
add_rect(slide, Inches(0.3), Inches(1.1), Inches(6.0), Inches(5.6), fill_color=C_PANEL)
add_rect(slide, Inches(0.3), Inches(1.1), Inches(0.05), Inches(5.6), fill_color=C_ACCENT)
add_textbox(slide, "⚠️  Keterbatasan & Analisis Kegagalan",
    left=Inches(0.5), top=Inches(1.2), width=Inches(5.7), height=Inches(0.45),
    font_size=14, bold=True, color=C_ACCENT2)

limitations = [
    ("❌ Class Imbalance", "BEBAS recall = 0,00. Weighted voting\nprediksi mayoritas (PENJARA_LAIN)."),
    ("❌ Ekstraksi Vonis 17,1%", "Hanya 6/35 vonis berhasil diekstrak.\nVariasi formulasi putusan sangat beragam."),
    ("❌ Dataset Kecil (35 dok.)", "28 training cases tidak cukup untuk\nmewakili variasi kasus pemalsuan."),
    ("⚙️ Revise & Retain", "Belum diimplementasikan sesuai\nsiklus CBR klasik Aamodt & Plaza."),
    ("📉 Retrieval MRR=0,0667", "Ground truth berbasis pasal terlalu\nsimplistis; relevansi lebih kompleks."),
]
for i, (title, desc) in enumerate(limitations):
    y = Inches(1.75) + i * Inches(0.92)
    add_textbox(slide, title,
        left=Inches(0.5), top=y, width=Inches(5.7), height=Inches(0.38),
        font_size=11, bold=True, color=C_WHITE)
    add_textbox(slide, desc,
        left=Inches(0.5), top=y + Inches(0.38), width=Inches(5.7), height=Inches(0.5),
        font_size=10, color=C_TEXT_DIM)

# Right column: conclusions + recommendations
add_rect(slide, Inches(6.65), Inches(1.1), Inches(6.35), Inches(5.6), fill_color=C_PANEL)
add_rect(slide, Inches(6.65), Inches(1.1), Inches(0.05), Inches(5.6), fill_color=C_GOLD)

add_textbox(slide, "✅  Kesimpulan",
    left=Inches(6.85), top=Inches(1.2), width=Inches(6.0), height=Inches(0.4),
    font_size=14, bold=True, color=C_GOLD)
add_multi_line_textbox(slide, [
    "• Sistem CBR berhasil diimplementasikan end-to-end",
    "• Akurasi prediksi 71,43% — performa moderat",
    "• Recall 100% untuk PENJARA_LAIN",
    "• MRR@5 = 0,0667 — perlu peningkatan retrieval",
],
    left=Inches(6.85), top=Inches(1.65), width=Inches(6.0), height=Inches(1.6),
    font_size=11, color=C_LIGHT_GRAY)

add_textbox(slide, "🔮  Rekomendasi Penelitian Lanjutan",
    left=Inches(6.85), top=Inches(3.35), width=Inches(6.0), height=Inches(0.4),
    font_size=14, bold=True, color=C_GOLD)
add_multi_line_textbox(slide, [
    "1. Augmentasi dataset → ≥100 dokumen putusan",
    "2. Fine-tuning IndoBERT pada korpus hukum ID",
    "3. Hybrid retrieval: BM25 + IndoBERT embedding",
    "4. Constraint-based filter (kecocokan pasal)",
    "5. Implementasi tahap Revise & Retain",
    "6. Confidence threshold untuk prediksi tidak pasti",
],
    left=Inches(6.85), top=Inches(3.8), width=Inches(6.0), height=Inches(2.6),
    font_size=11, color=C_LIGHT_GRAY)

add_textbox(slide, "05 / 05",
    left=W - Inches(1.3), top=H - Inches(0.4), width=Inches(1.0), height=Inches(0.3),
    font_size=10, color=C_TEXT_DIM, align=PP_ALIGN.RIGHT)

# ==============================================================
# SAVE
# ==============================================================
output_path = os.path.join(os.path.dirname(__file__), "presentasi_cbr.pptx")
prs.save(output_path)
print(f"[OK] File berhasil disimpan: {output_path}")
